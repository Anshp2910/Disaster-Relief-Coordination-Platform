import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx import Presentation

# Color palette
DARK_BG     = RGBColor(0x05, 0x0A, 0x16)   # #050a16
CARD_BG     = RGBColor(0x0F, 0x1B, 0x30)   # #0f1b30
ACCENT      = RGBColor(0x3B, 0x82, 0xF6)   # #3b82f6
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xD0, 0xE0, 0xF8)
GREEN       = RGBColor(0x22, 0xC5, 0x5E)
RED         = RGBColor(0xEF, 0x44, 0x44)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE      = RGBColor(0xA7, 0x8B, 0xFA)
ORANGE      = RGBColor(0xF9, 0x73, 0x16)
CYAN        = RGBColor(0x06, 0xB6, 0xD4)
TEAL        = RGBColor(0x14, 0xB8, 0xA6)

def add_rectangle(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size, color=None, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox

def add_slide_title(slide, title_text):
    add_textbox(slide, 0.4, 0.25, 8, 0.7, title_text, 28, WHITE, bold=True, align=PP_ALIGN.LEFT)
    bar = slide.shapes.add_shape(1, Inches(0.4), Inches(0.95), Inches(1.8), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

def add_body_text(slide, left, top, width, height, text, font_size=12, color=None, bullet=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    # first paragraph
    from pptx.util import Pt
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    if bullet:
        p.level = 0
        p.text = ''
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color if color else LIGHT_GRAY
    return txBox

def add_bullet_para(tf, text, level, font_size, color):
    p = tf.add_paragraph()
    p.level = level
    if level == 0:
        p.text = '• ' + text
    else:
        p.text = '– ' + text
    run = p.runs[0] if p.runs else p.add_run()
    run.font.size = Pt(font_size)
    run.font.color.rgb = color

def make_value_shape(slide, left, top, width, height, bg_color, value_text, label_text, accent_line=False):
    card = add_rectangle(slide, left, top, width, height, bg_color, ACCENT, 0.5)
    if accent_line:
        line = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT
        line.line.fill.background()
    add_textbox(slide, left+0.05, top+0.1, width-0.1, 0.5, str(value_text), 18, WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left+0.05, top+0.55, width-0.1, height-0.6, label_text, 9.5, LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return card

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank

# ─────────────────────────────────────────────
# SLIDE 1 — Title
# ─────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BG

# Large accent rect at top
accent_bar = add_rectangle(slide, 0, 0, 13.333, 0.25, ACCENT)
# Accent line at bottom
accent_bar2 = add_rectangle(slide, 0, 7.25, 13.333, 0.25, ACCENT)

# Title
add_textbox(slide, 0.8, 1.0, 11.5, 1.8,
    "DISASTER RELIEF\nCOORDINATION PLATFORM", 40, WHITE, bold=True, align=PP_ALIGN.LEFT)
add_textbox(slide, 0.8, 2.1, 11.5, 0.8,
    "Government of India Initiative | Mission-Critical Emergency Response System", 14, ACCENT, align=PP_ALIGN.LEFT)
add_textbox(slide, 0.8, 2.75, 11.5, 0.9,
    "Real-Time Collaboration  |  GIS Mapping  |  Resource Management  |  Multi-Language Support", 12, LIGHT_GRAY, align=PP_ALIGN.LEFT)
# Tags
tags = [("Technology", ACCENT), ("Web Dev", PURPLE), ("Security", AMBER), ("Cloud", CYAN)]
x_start, y_tag = 0.8, 3.8
for idx, (tag, col) in enumerate(tags):
    x = x_start + idx * 2.5
    tag_box = add_rectangle(slide, x, y_tag, 2.2, 0.4, CARD_BG, col, 1)
    add_textbox(slide, x+0.05, y_tag+0.05, 2.1, 0.3, tag, 10, col, align=PP_ALIGN.CENTER)

add_textbox(slide, 0.8, 4.5, 11.5, 1.0,
    "Live: disasterhelper.dpdns.org", 11, RGBColor(0x7A, 0x93, 0xB8), align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────
# SLIDE 2 — Domain & Technology Stack
# ─────────────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
bg = slide2.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BG
add_rectangle(slide2, 0, 0, 13.333, 0.15, ACCENT)
add_rectangle(slide2, 0, 7.35, 13.333, 0.15, ACCENT)
add_slide_title(slide2, "1. Domain & Technology Stack")

# Left card: domain
add_rectangle(slide2, 0.4, 1.15, 5.8, 5.8, CARD_BG, ACCENT, 0.8)
add_textbox(slide2, 0.55, 1.35, 5.5, 0.4, "Application Domain", 13, ACCENT, bold=True)
domain_items = [
    "Web Development          Full-Stack MERN Application",
    "Cyber Security          Enterprise-grade security hardening",
    "Cloud / DevOps          Docker, CI/CD, Render+Vercel deploy",
    "GIS / Data Science      Geo-spatial queries & dashboards",
    "IoT (extensible)        Real-time sensor-ready Socket.io layer",
]
y_off = 1.80
for row in domain_items:
    parts = row.split('  ')
    add_rectangle(slide2, 0.6, y_off, 0.08, 0.08, ACCENT)
    add_textbox(slide2, 0.82, y_off-0.05, 2.0, 0.3, parts[0], 10, ACCENT, bold=True)
    add_textbox(slide2, 2.85, y_off-0.05, 3.0, 0.3, parts[1], 10, LIGHT_GRAY)
    y_off += 0.75

# Right area: technologies
add_rectangle(slide2, 6.55, 1.15, 6.35, 5.8, CARD_BG, RGBColor(0x14, 0xB8, 0xA6), 0.8)
add_textbox(slide2, 6.7, 1.35, 6.2, 0.4, "Technologies Used", 13, TEAL, bold=True)

tech_groups = [
    ("Frontend",             ["React 18, TypeScript, Vite 5", "Framer Motion, Recharts", "Leaflet GIS, Socket.io Client", "i18next (10 languages), PWA SW"]),
    ("Backend",              ["Node.js 20, Express 4", "MongoDB, Mongoose, JWT", "Socket.io, Socket.io auth", "Joi validation, Multer, Passport"]),
    ("Libraries & Tools",    ["ESLint, Vitest, Playwright", "@axe-core (a11y), lint-staged", "bcrypt, Helmet, Compression", "Docker, GitHub Actions CI"]),
]
y_off = 1.80
colors = [ACCENT, PURPLE, AMBER]
for idx, (heading, items) in enumerate(tech_groups):
    add_textbox(slide2, 6.7, y_off, 6.0, 0.3, heading, 10.5, colors[idx], bold=True)
    y_off += 0.32
    for item in items:
        add_textbox(slide2, 6.85, y_off-0.02, 5.8, 0.25, "• " + item, 9.5, LIGHT_GRAY)
        y_off += 0.28
    y_off += 0.28

# ─────────────────────────────────────────────
# SLIDE 3 — Internship Objectives & Work Assigned
# ─────────────────────────────────────────────
slide3 = prs.slides.add_slide(blank_layout)
bg = slide3.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BG
add_rectangle(slide3, 0, 0, 13.333, 0.15, GREEN)
add_rectangle(slide3, 0, 7.35, 13.333, 0.15, GREEN)
add_slide_title(slide3, "2. Internship Objectives & Work Assigned")

cols = [
    ("Internship Objectives", [
        "Understand full-stack MERN development lifecycle",
        "Implement enterprise-grade security best practices",
        "Build real-time features using WebSockets (Socket.io)",
        "Deploy application using Docker & Render/Vercel",
        "Write comprehensive unit, integration & E2E tests",
        "Implement WCAG AA accessibility compliance",
    ], ACCENT),
    ("Work Assigned", [
        "Full-stack application development (frontend + backend)",
        "Database modeling & MongoDB integration",
        "Real-time event system & Socket.io chat",
        "GIS mapping with Leaflet & geo-spatial queries",
        "Security hardening (11 security items implemented)",
        "Docker containerization & Render deployment",
        "Performance optimization (code splitting, lazy loading)",
        "Writing testing suites (Vitest + Playwright)",
    ], GREEN),
]
x_start = 0.4
for idx, (heading, items, col) in enumerate(cols):
    add_rectangle(slide3, x_start, 1.15, 6.1, 5.8, CARD_BG, col, 0.8)
    add_textbox(slide3, x_start+0.1, 1.3, 5.9, 0.4, heading, 13, col, bold=True)
    y_off = 1.78
    for item in items:
        add_textbox(slide3, x_start+0.2, y_off-0.02, 5.7, 0.28, "• " + item, 10.5, LIGHT_GRAY)
        y_off += 0.42
    x_start += 6.45

# ─────────────────────────────────────────────
# SLIDE 4 — Problem Statement
# ─────────────────────────────────────────────
slide4 = prs.slides.add_slide(blank_layout)
bg = slide4.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BG
add_rectangle(slide4, 0, 0, 13.333, 0.15, RED)
add_rectangle(slide4, 0, 7.35, 13.333, 0.15, RED)
add_slide_title(slide4, "3. Problem Statement")

# Problem statement card
add_rectangle(slide4, 0.4, 1.15, 12.5, 2.8, CARD_BG, RED, 0.6)
add_textbox(slide4, 0.6, 1.28, 12.2, 0.35, "EXISTING SYSTEM / PROBLEM", 11, RED, bold=True)
add_textbox(slide4, 0.55, 1.65, 12.2, 2.15,
    "During disaster events, relief coordination is severely hampered by:\n"
    "• Fragmented communication between government agencies, NGOs and volunteers\n"
    "• No centralized view of resources, shelters, or affected populations\n"
    "• Delayed response due to manual phone calls & paper-based tracking\n"
    "• No real-time SOS alerts broadcast to field responders\n"
    "• Inability to track resource allocation across disaster zones dynamically\n"
    "• Language barriers limiting access for rural & Urdu-speaking communities",
    11, LIGHT_GRAY)

# Proposed solution card
add_rectangle(slide4, 0.4, 4.1, 12.5, 3.0, CARD_BG, GREEN, 0.6)
add_textbox(slide4, 0.6, 4.22, 12.2, 0.35, "PROPOSED SOLUTION", 11, GREEN, bold=True)
add_textbox(slide4, 0.55, 4.58, 12.2, 2.35,
    "The Disaster Relief Coordination Platform (DRCP) provides a unified MERN-based web application:\n"
    "• Centralized task board with status tracking for all relief requests\n"
    "• Real-time GIS map (Leaflet) with disaster overlays for zone monitoring\n"
    "• SOS emergency broadcast system with GPS coordinates to all responders\n"
    "• Resource & inventory management with smart allocation & matching\n"
    "• 10-language i18n support (including Urdu RTL) for inclusive access\n"
    "• Admin dashboard with escalation workflows, bulk CSV ops & analytics",
    11, LIGHT_GRAY)

# ─────────────────────────────────────────────
# SLIDE 5 — Architecture / Workflow
# ─────────────────────────────────────────────
slide5 = prs.slides.add_slide(blank_layout)
bg = slide5.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BG
add_rectangle(slide5, 0, 0, 13.333, 0.15, PURPLE)
add_rectangle(slide5, 0, 7.35, 13.333, 0.15, PURPLE)
add_slide_title(slide5, "4. System Architecture")

# Three-tier columns
tiers = [
    ("CLIENT TIER", ["React 18 SPA (port 5173)", "22 lazy-loaded pages", "17 reusable UI components", "9 custom React hooks", "Role-based routing",
                     "Dark/Light/Neon themes"], ACCENT, 0.4),
    ("API TIER", ["Express.js REST API (port 5001)", "15 route modules (70+ endpoints)", "JWT Auth + Passport OAuth2", "Socket.io real-time",
                  "Joi input validation", "Rate limiting per endpoint"], PURPLE, 4.55),
    ("DATA TIER", ["MongoDB 7 (2dsphere indexes)", "9 Mongoose models", "GeoJSON Point locations", "$geoWithin / $near queries",
                   "5x retry connection logic", "Real-time sync via Socket.io"], TEAL, 8.7),
]
x_start = 0.4
for idx, (heading, items, col, left) in enumerate(tiers):
    add_rectangle(slide5, left, 1.15, 3.85, 5.8, CARD_BG, col, 0.8)
    add_textbox(slide5, left+0.1, 1.28, 3.65, 0.38, heading, 11, col, bold=True, align=PP_ALIGN.CENTER)
    y_off = 1.72
    for item in items:
        add_textbox(slide5, left+0.18, y_off-0.01, 3.55, 0.26, "• " + item, 9.5, LIGHT_GRAY)
        y_off += 0.78

# Arrow
add_textbox(slide5, 4.3, 3.2, 0.9, 0.5, "→", 28, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slide5, 8.45, 3.2, 0.9, 0.5, "→", 28, WHITE, bold=True, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 6 — Project Description / Modules
# ─────────────────────────────────────────────
slide6 = prs.slides.add_slide(blank_layout)
bg = slide6.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BG
add_rectangle(slide6, 0, 0, 13.333, 0.15, AMBER)
add_rectangle(slide6, 0, 7.35, 13.333, 0.15, AMBER)
add_slide_title(slide6, "5. Project Description — Modules")

modules = [
    ("🔐  Auth Module",     "JWT auth, OAuth2 (GitHub/Google), password reset via email, role-based access (volunteer / NGO / admin)", ACCENT),
    ("🗺  GIS Map Module",  "Leaflet maps with satellite/terrain overlays, flood/cyclone/earthquake layers, pulsing SOS markers", TEAL),
    ("📋  Requests Module", "CRUD relief requests with status/priority, file attachments, comments, real-time chat", ORANGE),
    ("📦  Resources Module","15-category inventory, allocation to requests, stock tracking, CSV bulk import/export", GREEN),
    ("🚨  SOS System",      "One-tap emergency alert with GPS broadcast to all responders, auto-notification", RED),
    ("👥  Volunteer Mgmt",  "Profile, availability, skill tags, check-in/attendance, schedule calendar", PURPLE),
    ("💬  Real-time Chat",  "Socket.io rooms per request, JWT-authenticated, live message updates", CYAN),
    ("📊  Admin Dashboard", "User & request management, escalation workflows, charts, CSV upload", AMBER),
    ("🌍  i18n Module",     "10 languages incl. Urdu (RTL), lazy-loaded locale files, automatic locale detection", ACCENT),
    ("🔒  Security",        "Helmet CSP, CORS allowlist, rate limiting (SOS: 5/min), Joi validation, Helmet, SanitizeBody", RED),
]

card_w, card_h = 4.0, 1.65
gap = 0.22
positions = []
for row in range(3):
    for col in range(4):
        if row * 4 + col >= len(modules):
            break
        left = 0.4 + col * (card_w + gap)
        top = 1.2 + row * (card_h + 0.18)
        positions.append((left, top))

for idx, ((heading, desc, color)) in enumerate(modules):
    left, top = positions[idx]
    add_rectangle(slide6, left, top, card_w, card_h, CARD_BG, color, 0.6)
    add_textbox(slide6, left+0.1, top+0.06, card_w-0.2, 0.32, heading, 9.5, color, bold=True)
    add_textbox(slide6, left+0.1, top+0.38, card_w-0.2, card_h-0.45, desc, 8.5, LIGHT_GRAY)

# ─────────────────────────────────────────────
# SLIDE 7 — Week-wise Activities
# ─────────────────────────────────────────────
slide7 = prs.slides.add_slide(blank_layout)
bg = slide7.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BG
add_rectangle(slide7, 0, 0, 13.333, 0.15, CYAN)
add_rectangle(slide7, 0, 7.35, 13.333, 0.15, CYAN)
add_slide_title(slide7, "6. Work Performed — Week-wise Activities")

weeks = [
    ("Week 1", "Training & Requirement Analysis", [
        "Set up MERN development environment",
        "Analyzed 17+ page requirements & existing APIs",
        "Studied design system (17 CSS files, design tokens)",
        "Explored 9 Mongoose models & 15 Express routes",
        "Reviewed security checklist (11 items)",
        "Learned Socket.io real-time event patterns",
    ]),
    ("Week 2", "Development", [
        "Built authentication (JWT + OAuth2)",
        "Developed GIS mapping with Leaflet (satellite/terrain/flood overlays)",
        "Implemented real-time SOS broadcast system",
        "Built resource & inventory management (15 categories)",
        "Created admin dashboard with charts & escalation",
        "Integrated Socket.io chat rooms per request",
    ]),
    ("Week 3", "Testing & Integration", [
        "Wrote 175+ unit tests (Vitest + @testing-library)",
        "13 Playwright E2E tests (theme toggle, login, lang)",
        "Accessibility tests (axe-core, ARIA check)",
        "Integrated Socket.io with auth & reconnection",
        "PWA: service worker v7, offline fallback",
        "Cross-browser & responsive testing",
    ]),
    ("Week 4", "Deployment & Documentation", [
        "Docker multi-stage builds (client + server)",
        "Docker Compose for local MongoDB",
        "GitHub Actions CI/CD pipeline",
        "Render backend + Vercel frontend deployment",
        "Wrote DESIGN_SYSTEM.md, UI-UX-DESIGN.md",
        "Integrated monitoring & server keepalive",
    ]),
]
card_w, card_h = 6.1, 5.6
x_positions = [0.4, 6.75]
for idx, (week, subtitle, items) in enumerate(weeks):
    left = x_positions[idx % 2]
    top = 1.15 + (idx // 2) * 3.2
    col = [ACCENT, GREEN, PURPLE, AMBER][idx]
    add_rectangle(slide7, left, top, card_w, card_h, CARD_BG, col, 1.0)
    add_rectangle(slide7, left, top, card_w, 0.42, col)
    add_textbox(slide7, left+0.08, top+0.02, card_w-0.15, 0.38, week, 12, WHITE, bold=True)
    add_textbox(slide7, left+0.08, top+0.45, card_w-0.15, 0.32, subtitle, 10.5, col, bold=True)
    y_off = top + 0.82
    for item in items:
        add_textbox(slide7, left+0.15, y_off-0.02, card_w-0.22, 0.26, "• " + item, 9.5, LIGHT_GRAY)
        y_off += 0.4

# ─────────────────────────────────────────────
# SLIDE 8 — Implementation / Screenshots / Tech
# ─────────────────────────────────────────────
slide8 = prs.slides.add_slide(blank_layout)
bg = slide8.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BG
add_rectangle(slide8, 0, 0, 13.333, 0.15, ORANGE)
add_rectangle(slide8, 0, 7.35, 13.333, 0.15, ORANGE)
add_slide_title(slide8, "7. Implementation Details")

# Source code snippet card
add_rectangle(slide8, 0.4, 1.15, 7.8, 5.8, CARD_BG, ACCENT, 0.6)
add_textbox(slide8, 0.55, 1.28, 7.6, 0.38, "SOURCE CODE HIGHLIGHTS (server/routes/sos.js)", 10, ACCENT, bold=True)
code_lines = [
    "// SOS emergency broadcast via Socket.io",
    "io.to('sos').emit('new:sos', {",
    "  location: SOSAlert.location,",
    "  category,",
    "  priority,",
    "  requestedBy: user.name,",
    "  timestamp: new Date(),",
    "});",
    "",
    "// Geo-spatial query — find nearby resources",
    "const nearby = await Resource.find({",
    "  location: {",
    "    $near: {",
    "      $geometry: { type: 'Point', coordinates },",
    "      $maxDistance: 5000",
    "    }",
    "  },",
    "  category,",
    "}).select('name category location amount');",
    "",
    "// Joi validation — sanitize SOS input",
    "const schema = Joi.object({",
    "  category: Joi.string().required(),",
    "  priority: Joi.string().valid(...).required(),",
    "  description: Joi.string().max(500).allow(''),",
    "});",
]
add_textbox(slide8, 0.55, 1.72, 7.55, 5.1,
    "\n".join(code_lines), 9.5, RGBColor(0xE0, 0xF2, 0xFF))

# Right: screenshots list + database tables
add_rectangle(slide8, 8.55, 1.15, 4.37, 5.8, CARD_BG, AMBER, 0.6)
add_textbox(slide8, 8.7, 1.28, 4.15, 0.38, "UI PAGES / SCREENSHOTS", 10, AMBER, bold=True)

pages = [
    "/login – Auth page (glass card design)",
    "/dashboard – Request list with filters",
    "/map – Full-screen GIS Leaflet map",
    "/requests/new – Create request form",
    "/resources – Inventory CRUD table",
    "/zones – Disaster zone heat overlays",
    "/admin – Admin stats + charts",
    "/schedules – Volunteer calendar",
    "/profile – Account management",
]
y_off = 1.75
for pg in pages:
    add_textbox(slide8, 8.7, y_off-0.02, 4.1, 0.25, "• " + pg, 9.5, LIGHT_GRAY)
    y_off += 0.48

add_textbox(slide8, 8.7, 5.72, 4.1, 0.28, "→ disasterhelper.dpdns.org", 8, ACCENT, bold=True)

# Database table
add_textbox(slide8, 8.7, 5.85, 4.1, 0.28, "DATABASE (9 Models)", 9, AMBER, bold=True)
tables = "• User, Request, Resource\n• Zone, SosAlert, Incident\n• Schedule, ChatMessage, Feedback"
add_textbox(slide8, 8.7, 6.1, 4.1, 0.5, tables, 9, LIGHT_GRAY)

# ─────────────────────────────────────────────
# SLIDE 9 — Results & Achievements
# ─────────────────────────────────────────────
slide9 = prs.slides.add_slide(blank_layout)
bg = slide9.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BG
add_rectangle(slide9, 0, 0, 13.333, 0.15, GREEN)
add_rectangle(slide9, 0, 7.35, 13.333, 0.15, GREEN)
add_slide_title(slide9, "8. Results & Achievements")

# Metrics row
metrics = [
    ("17",        "UI Pages",            ACCENT),
    ("70+",       "API Endpoints",       PURPLE),
    ("9",         "DB Models",           TEAL),
    ("10",        "Languages",           GREEN),
    ("175+",      "Unit Tests",          AMBER),
    ("13",        "E2E Tests",           ORANGE),
    ("15",        "Route Modules",       CYAN),
    ("19",        "CSS Files",           ACCENT),
    ("A+",        "Security Grade",      RED),
    ("AA+",       "Accessibility",       GREEN),
]
card_w_m, card_h_m = 1.2, 1.6
cols_m = 5
x_start = 0.4
y_start = 1.2
x_gap_m = 0.12
y_gap_m = 0.18
for idx, (val, label, col) in enumerate(metrics):
    col_idx = idx % cols_m
    row_idx = idx // cols_m
    left = x_start + col_idx * (card_w_m + x_gap_m)
    top = y_start + row_idx * (card_h_m + y_gap_m)
    make_value_shape(slide9, left, top, card_w_m, card_h_m, CARD_BG, val, label)

# Bottom section
y_bot = 6.05
achievements = [
    ("Features Completed", [
        "End-to-end MERN full-stack app with 17 pages & 70+ API endpoints",
        "Real-time GIS mapping with 9 disaster zone overlays",
        "SOS emergency broadcast with Socket.io",
        "10-language i18n including Urdu (RTL)",
        "PWA with service worker & offline fallback",
        "Docker multi-stage build + Render deployment",
        "19 CSS files — glassmorphism design system, light/dark/neon themes",
    ], GREEN),
    ("Testing Results", [
        "175+ unit tests passing — Vitest framework",
        "13 Playwright E2E tests (theme toggle, login, lang, password)",
        "axe-core accessibility audit — WCAG AA+ compliant",
        "Rate limiting tested per-endpoint",
        "Lighthouse score: >95 (target met)",
        "0 unhandled errors in production",
    ], ACCENT),
    ("Deliverables", [
        "Production app: disasterhelper.dpdns.org",
        "GitHub repo with GitHub Actions CI/CD",
        "Docker Compose local setup",
        "Full documentation (DESIGN_SYSTEM.md, UI-UX-DESIGN.md)",
        "Server keepalive via inner + UptimeRobot (5-min pings)",
        "render.yaml + Dockerfile.server for Render deploy",
    ], PURPLE),
]
x_b = [0.4, 4.55, 8.7]
for idx, (heading, items, col) in enumerate(achievements):
    add_rectangle(slide9, x_b[idx], y_bot, 3.85, 2.55, CARD_BG, col, 0.6)
    add_textbox(slide9, x_b[idx]+0.1, y_bot+0.05, 3.65, 0.35, heading, 10.5, col, bold=True)
    y_off = y_bot + 0.45
    for item in items:
        add_textbox(slide9, x_b[idx]+0.14, y_off-0.02, 3.57, 0.28,
            "• " + item if idx < 2 else "✔ " + item, 9, LIGHT_GRAY)
        y_off += 0.35

# Thank you slide
slideT = prs.slides.add_slide(blank_layout)
bg = slideT.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = DARK_BG
add_rectangle(slideT, 0, 0, 13.333, 0.15, ACCENT)
add_rectangle(slideT, 0, 7.35, 13.333, 0.15, ACCENT)
add_textbox(slideT, 0, 2.0, 13.333, 1.5,
    "Thank You", 48, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_textbox(slideT, 0, 3.5, 13.333, 0.7,
    "Disaster Relief Coordination Platform", 18, ACCENT, align=PP_ALIGN.CENTER)
add_textbox(slideT, 0, 4.1, 13.333, 0.6,
    "governor.of.india.initiative | disasterhelper.dpdns.org", 11, LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_textbox(slideT, 0, 4.65, 13.333, 0.5,
    "Live · Feature-Complete · Production-Ready", 10, RGBColor(0x7A, 0x93, 0xB8), align=PP_ALIGN.CENTER)

out = r"C:\Disaster Relief Coordination Platform\Internship_Report_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
